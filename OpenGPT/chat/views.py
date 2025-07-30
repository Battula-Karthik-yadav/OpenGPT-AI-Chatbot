import requests
import fitz
from docx import Document
import pytesseract
import json
from PIL import Image
from io import BytesIO
from datetime import datetime
from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse, FileResponse, StreamingHttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.forms import AuthenticationForm, UserCreationForm
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.core.files.storage import default_storage
from django.utils.timezone import now

from reportlab.pdfgen import canvas
from pptx import Presentation

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

from .models import ChatSession, ChatMessage, UploadedFile, GeneratedFile

# ========== BASIC VIEWS ==========

def home(request):
    return render(request, 'home.html')

def logout_view(request):
    logout(request)
    return redirect('home')

def custom_login_view(request):
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            login(request, form.get_user())
            return redirect('chat')
        else:
            messages.error(request, "Invalid login credentials")
    else:
        form = AuthenticationForm()
    return render(request, 'login.html', {'form': form})

def custom_register_view(request):
    if request.method == 'POST':
        form = UserCreationForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('login')
    else:
        form = UserCreationForm()
    return render(request, 'register.html', {'form': form})

# ========== CHAT VIEWS ==========

@login_required
@login_required
def chat_view(request):
    # Create new chat session
    new_session = ChatSession.objects.create(user=request.user, title="New Chat")

    # Fetch all previous sessions (except the new one)
    previous_sessions = ChatSession.objects.filter(
        user=request.user, is_deleted=False
    ).exclude(id=new_session.id).order_by('-updated_at')

    #remberCD
    # Group sessions by "Today", "Yesterday", "April 25, 2025", etc.
    grouped_sessions = {}
    for session in previous_sessions:
        group_label = session.get_date_group()
        grouped_sessions.setdefault(group_label, []).append(session)

    return render(request, 'chat.html', {
        'currentSessionId': new_session.id,
        'groupedSessions': grouped_sessions,  # <-- pass the grouped sessions instead of flat list
    })

@csrf_exempt
@require_POST
@login_required
def send_message(request):
    user = request.user
    message_text = request.POST.get("message", "").strip()
    session_id = request.POST.get("session_id")
    uploaded_files = request.FILES.getlist('file')

    if not message_text and not uploaded_files:
        return JsonResponse({"error": "Missing message or file"}, status=400)

    session = get_object_or_404(ChatSession, id=session_id, user=user, is_deleted=False)

    def stream_response():
        # Stream normal text message (if provided)
        if message_text:
            ChatMessage.objects.create(session=session, role='user', content=message_text)
            for chunk in stream_assistant_response(message_text, session):
                yield chunk
            yield "\n"  # Separate messages visually

        # Stream uploaded files one by one
        for uploaded_file in uploaded_files:
            file_content = extract_file_content(uploaded_file)
            if file_content is None:
                yield f"\n[Skipped file: {uploaded_file.name}]"
                continue

            file_message_content = f"[File Upload: {uploaded_file.name}]\n{file_content}"
            ChatMessage.objects.create(session=session, role='user', content=file_message_content)

            for chunk in stream_assistant_response(file_message_content, session):
                yield chunk
            yield "\n"  # Separate each file's assistant response

        session.save()

    return StreamingHttpResponse(stream_response(), content_type="text/plain")

@csrf_exempt
@require_POST
@login_required
def new_chat_session(request):
    session = ChatSession.objects.create(user=request.user, title="New Chat")
    return JsonResponse({"session_id": session.id})

def extract_file_content(uploaded_file):
    try:
        file_name = uploaded_file.name.lower()

        if file_name.endswith(".pdf"):
            pdf_doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            return "\n".join(page.get_text() for page in pdf_doc)

        elif file_name.endswith(".docx"):
            doc = Document(uploaded_file)
            return "\n".join(para.text for para in doc.paragraphs)

        elif file_name.endswith(('.png', '.jpg', '.jpeg')):
            image = Image.open(uploaded_file)
            text = pytesseract.image_to_string(image)
            return text if text.strip() else None

        elif file_name.endswith(('.txt', '.md', '.json')):
            return uploaded_file.read().decode('utf-8')

        else:
            return None
    except Exception as e:
        print(f"Error extracting file content: {e}")
        return None

def stream_assistant_response(message_content, session):
    try:
        response = requests.post(
            "http://localhost:11434/api/chat",
            json={
                "model": "mistral",
                "messages": [{"role": "user", "content": message_content}],
                "stream": True,  # Streaming from the assistant
                "timeout" : 60,
            },
            stream=True,

        )

        assistant_full_reply = ""

        for line in response.iter_lines():
            if line:
                try:
                    data = json.loads(line.decode('utf-8'))
                    content_piece = data.get("message", {}).get("content", "")
                    assistant_full_reply += content_piece
                    yield content_piece
                except json.JSONDecodeError:
                    continue

        # Save full assistant reply after streaming
        ChatMessage.objects.create(session=session, role='assistant', content=assistant_full_reply)

    except requests.exceptions.RequestException as e:
        yield f"\n[Assistant request error: {str(e)}]"

@login_required
def get_chat_history(request):
    sessions = ChatSession.objects.filter(user=request.user, is_deleted=False).order_by('-updated_at')
    chat_data = [{
        'id': s.id,
        'title': s.title,
        'created_at': s.created_at.strftime('%Y-%m-%d %H:%M:%S'),
        'updated_at': s.updated_at.strftime('%Y-%m-%d %H:%M:%S'),
    } for s in sessions]
    return JsonResponse({"sessions": chat_data})

@csrf_exempt
@login_required
def load_session_messages(request, session_id):
    session = get_object_or_404(ChatSession, id=session_id, user=request.user, is_deleted=False)
    messages = ChatMessage.objects.filter(session=session).order_by('created_at')

    messages_data = [{
        'role': msg.role,
        'content': msg.content,
        'timestamp': msg.created_at.strftime('%Y-%m-%d %H:%M:%S')
    } for msg in messages]

    return JsonResponse({'messages': messages_data})

@csrf_exempt
@require_POST
@login_required
def rename_chat(request):
    session_id = request.POST.get('session_id')
    new_title = request.POST.get('new_title')

    if not session_id or not new_title:
        return JsonResponse({"error": "Missing session ID or new title"}, status=400)

    session = get_object_or_404(ChatSession, id=session_id, user=request.user, is_deleted=False)
    session.rename(new_title)

    return JsonResponse({"success": True, "new_title": new_title})

@csrf_exempt
@require_POST
@login_required
def delete_chat(request):
    session_id = request.POST.get('session_id')

    if not session_id:
        return JsonResponse({"error": "Missing session ID"}, status=400)

    session = get_object_or_404(ChatSession, id=session_id, user=request.user)
    session.soft_delete()

    return JsonResponse({"success": True})

@csrf_exempt
@login_required
def search_chats(request):
    query = request.GET.get('q', '')
    sessions = ChatSession.objects.filter(user=request.user, is_deleted=False, title__icontains=query).order_by('-updated_at')

    search_data = [{
        'id': s.id,
        'title': s.title,
        'updated_at': s.updated_at.strftime('%Y-%m-%d %H:%M:%S')
    } for s in sessions]

    return JsonResponse({'results': search_data})

# ========== FILE UPLOAD VIEWS ==========

@csrf_exempt
@login_required
def upload_pdf(request):
    if request.method == 'POST' and request.FILES:
        uploaded_files = []
        for file in request.FILES.getlist('files'):
            filename = default_storage.save(file.name, file)
            uploaded_files.append({
                'name': file.name,
                'url': default_storage.url(filename)
            })
        return JsonResponse({'status': 'success', 'files': uploaded_files})
    return JsonResponse({'status': 'error', 'message': 'No files uploaded.'})


@csrf_exempt
@login_required
def upload_image(request):
    if request.method == 'POST' and request.FILES.get('image'):
        image_file = request.FILES['image']

        # Ensure the file is an image type
        if not image_file.content_type.startswith('image/'):
            return JsonResponse({"error": "Unsupported file type!"}, status=400)

        try:
            image = Image.open(image_file)
            image.thumbnail((100, 100))
            # Save or process the image as required
            return JsonResponse({"message": "Image uploaded and processed successfully!"})
        except Exception as e:
            return JsonResponse({"error": f"Error processing image: {str(e)}"}, status=400)
    return JsonResponse({"error": "No image file provided!"}, status=400)


@csrf_exempt
@require_POST
@login_required
def generate_pdf(request):
    content = request.POST.get('content', '').strip()
    if not content:
        return JsonResponse({"error": "No content provided"}, status=400)

    buffer = BytesIO()
    p = canvas.Canvas(buffer)
    p.drawString(100, 800, content[:4000])  # Ensure content is not too large
    p.showPage()
    p.save()
    buffer.seek(0)

    return FileResponse(buffer, as_attachment=True, filename="generated.pdf")

@csrf_exempt
@require_POST
@login_required
def generate_ppt(request):
    content = request.POST.get('content', '').strip()
    if not content:
        return JsonResponse({"error": "No content provided"}, status=400)

    prs = Presentation()
    for i, line in enumerate(content.split('\n')):
        if line.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = f"Slide {i+1}"
            body.text = line.strip()

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return FileResponse(buffer, as_attachment=True, filename="generated.pptx")
